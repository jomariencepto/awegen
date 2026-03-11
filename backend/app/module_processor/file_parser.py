import os
import re
import unicodedata

import PyPDF2
import docx
from docx.oxml.ns import qn
from pptx import Presentation
from openpyxl import load_workbook
from app.utils.logger import get_logger

logger = get_logger(__name__)

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
    logger.info("pdfplumber available — using as primary PDF extractor")
except ImportError:
    HAS_PDFPLUMBER = False
    logger.warning("pdfplumber not installed — PDF extraction will use PyPDF2 only")

try:
    import pytesseract
    HAS_PYTESSERACT = True
    # Try common Windows install path; ignored if Tesseract is already on PATH
    _TESSERACT_WIN = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    if os.path.exists(_TESSERACT_WIN):
        pytesseract.pytesseract.tesseract_cmd = _TESSERACT_WIN
    logger.info("pytesseract available — image OCR enabled")
except ImportError:
    HAS_PYTESSERACT = False
    logger.warning("pytesseract not installed — image OCR unavailable")

IMAGE_EXTENSIONS = {'jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp'}
MAX_FILE_MB = float(os.getenv("AI_MAX_FILE_MB", "25"))
MAX_PDF_PAGES = int(os.getenv("AI_MAX_PDF_PAGES", "500"))


class FileParser:
    """
    Comprehensive file parser for educational module documents.

    Extraction strategy:
    - PDF   : pdfplumber (primary, handles multi-column + tables) → PyPDF2 (fallback)
    - DOCX  : XML element-tree traversal in document order (paragraphs + tables interleaved)
              + section headers and footers
    - PPTX  : all shapes (recursive through groups) + speaker notes
    - XLSX  : all sheets, all rows
    - TXT   : multi-encoding fallback (UTF-8-sig → UTF-8 → Latin-1 → CP1252)

    All output passes through _clean_extracted_text() which normalises Unicode
    ligatures, curly quotes, non-breaking spaces and control characters.
    """

    # ------------------------------------------------------------------
    # Unicode normalisation map: ligatures, smart quotes, dashes, spaces
    # ------------------------------------------------------------------
    _LIGATURE_MAP = {
        '\ufb00': 'ff',   # ﬀ
        '\ufb01': 'fi',   # ﬁ
        '\ufb02': 'fl',   # ﬂ
        '\ufb03': 'ffi',  # ﬃ
        '\ufb04': 'ffl',  # ﬄ
        '\u2018': "'",    # left single quotation mark
        '\u2019': "'",    # right single quotation mark / apostrophe
        '\u201c': '"',    # left double quotation mark
        '\u201d': '"',    # right double quotation mark
        '\u2013': '-',    # en dash
        '\u2014': '-',    # em dash
        '\u2022': '-',    # bullet •
        '\u00a0': ' ',    # non-breaking space
        '\u00ad': '',     # soft hyphen (invisible; remove)
    }

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    @staticmethod
    def parse_file(file_path, file_type):
        """
        Parse a file and return its full text content.

        Args:
            file_path (str): Absolute path to the file.
            file_type (str): File extension without the dot (e.g. 'pdf', 'docx').

        Returns:
            str | None: Extracted and cleaned text, or None on failure.
        """
        try:
            # Resource guards
            size_mb = os.path.getsize(file_path) / (1024 * 1024)
            if size_mb > MAX_FILE_MB:
                logger.error(f"File too large ({size_mb:.1f}MB) > limit {MAX_FILE_MB}MB")
                raise ValueError("file_size_limit_exceeded")

            ft = file_type.lower()
            if ft == 'pdf':
                try:
                    with open(file_path, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        if len(reader.pages) > MAX_PDF_PAGES:
                            logger.error(f"PDF page count {len(reader.pages)} exceeds limit {MAX_PDF_PAGES}")
                            raise ValueError("pdf_page_limit_exceeded")
                except ValueError:
                    raise
                except Exception as count_err:
                    logger.warning(f"Could not count PDF pages: {count_err}")
                return FileParser._parse_pdf(file_path)
            elif ft in ('doc', 'docx'):
                return FileParser._parse_docx(file_path)
            elif ft in ('ppt', 'pptx'):
                return FileParser._parse_pptx(file_path)
            elif ft in ('xls', 'xlsx'):
                return FileParser._parse_xlsx(file_path)
            elif ft == 'txt':
                return FileParser._parse_txt(file_path)
            elif ft in IMAGE_EXTENSIONS:
                return FileParser._parse_image(file_path)
            else:
                logger.error(f"Unsupported file type: {file_type}")
                return None
        except Exception as e:
            logger.error(f"Error parsing file {file_path}: {str(e)}", exc_info=True)
            return None

    # ------------------------------------------------------------------
    # Shared text cleaner
    # ------------------------------------------------------------------

    @staticmethod
    def _fix_char_per_line(text):
        """
        Repair PDF extraction artifact where each glyph is stored on its own line.

        Type 1 — newline-per-character (most common in scanned/positioned PDFs):
            "C\\no\\nm\\np\\nl\\ne\\nt\\ne\\n\\nt\\nh\\ne" → "Complete the"
            Word breaks are signalled by empty lines between single-char runs.
            When no empty-line separators exist, wordninja is used to segment
            the concatenated characters into proper English words.

        Type 2 — space-per-character:
            "C o m p l e t e" → "Complete"

        Processed paragraph-by-paragraph so normal paragraphs are left intact.
        """
        try:
            import wordninja as _wn
            _wordninja_available = True
        except ImportError:
            _wn = None
            _wordninja_available = False

        def _segment_if_squished(token):
            if _wordninja_available and len(token) > 20:
                parts = _wn.split(token)
                if len(parts) > 1:
                    return parts
            return [token]

        paragraphs = re.split(r'\n{2,}', text)
        fixed = []

        for para in paragraphs:
            lines = para.split('\n')
            total = len(lines)
            if total == 0:
                fixed.append(para)
                continue

            single_char = sum(1 for l in lines if len(l.strip()) == 1)

            if total >= 2 and single_char / total > 0.45:
                # Type 1: each letter on its own line; empty lines = word breaks
                words = []
                current_chars = []
                for line in lines:
                    c = line.strip()
                    if len(c) == 1:
                        current_chars.append(c)
                    else:
                        if current_chars:
                            joined = ''.join(current_chars)
                            # Skip multi-char line if it is just the squished copy
                            c_alpha = re.sub(r'[^A-Za-z0-9]', '', c).lower()
                            joined_alpha = re.sub(r'[^A-Za-z0-9]', '', joined).lower()
                            is_squished_duplicate = (
                                c_alpha == joined_alpha
                                or (len(c_alpha) > 10 and joined_alpha.startswith(c_alpha[:10]))
                            )
                            words.extend(_segment_if_squished(joined))
                            current_chars = []
                            if c and not is_squished_duplicate:
                                words.extend(_segment_if_squished(c))
                        elif c:
                            words.extend(_segment_if_squished(c))
                if current_chars:
                    joined = ''.join(current_chars)
                    words.extend(_segment_if_squished(joined))
                fixed.append(' '.join(words))
            else:
                # Type 2: space-per-character — join runs of 4+ single-letter tokens
                fixed_para = re.sub(
                    r'(?<!\w)([A-Za-z] ){4,}[A-Za-z](?!\w)',
                    lambda m: m.group(0).replace(' ', ''),
                    para
                )
                fixed.append(fixed_para)

        return '\n\n'.join(fixed)

    @staticmethod
    def _clean_extracted_text(text):
        """
        Normalise Unicode artifacts, ligatures, curly quotes, and whitespace.

        Operations (in order):
        1. Replace known ligatures and special characters.
        2. NFKC Unicode normalisation (decomposes compatibility characters).
        3. Strip C0/C1 control characters (keep \\n and \\t).
        4. Fix character-per-line PDF extraction artifact.
        5. Collapse runs of 4+ blank lines down to 2.
        6. Strip leading/trailing whitespace.
        """
        if not text:
            return text

        # 1. Replace known ligatures / smart punctuation
        for char, replacement in FileParser._LIGATURE_MAP.items():
            text = text.replace(char, replacement)

        # 2. Unicode NFKC normalisation
        text = unicodedata.normalize('NFKC', text)

        # 3. Remove control characters except newline (\n) and tab (\t)
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)

        # 4. Fix character-per-line / character-per-space PDF artifacts
        text = FileParser._fix_char_per_line(text)

        # 5. Collapse excessive blank lines
        text = re.sub(r'\n{4,}', '\n\n\n', text)

        return text.strip()

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------

    @staticmethod
    def _parse_pdf(file_path):
        """
        Extract text from a PDF file.

        Primary  : pdfplumber — handles multi-column layouts and structured tables.
                   Tables are emitted as tab-separated rows immediately after body text.
        Fallback : PyPDF2 — fast text-layer extraction, single-column only.
        """
        text = ''

        # --- Primary: pdfplumber -----------------------------------------------
        if HAS_PDFPLUMBER:
            try:
                page_texts = []
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        parts = []

                        # Body text (multi-column aware)
                        body = page.extract_text(x_tolerance=2, y_tolerance=2)
                        if body:
                            parts.append(body)

                        # Structured tables → tab-separated rows
                        for table in page.extract_tables():
                            for row in table:
                                row_str = '\t'.join(
                                    (cell or '').strip() for cell in row
                                )
                                if row_str.strip():
                                    parts.append(row_str)

                        if parts:
                            page_texts.append('\n'.join(parts))

                text = '\n\n'.join(page_texts)

                if text.strip():
                    logger.info(f"pdfplumber extracted {len(text):,} chars from {file_path}")
                    return FileParser._clean_extracted_text(text)

                logger.warning("pdfplumber returned empty text — falling back to PyPDF2")

            except Exception as e:
                logger.warning(f"pdfplumber failed ({e}) — falling back to PyPDF2")

        # --- Fallback: PyPDF2 --------------------------------------------------
        try:
            parts = []
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        parts.append(page_text)
            text = '\n'.join(parts)
            logger.info(f"PyPDF2 extracted {len(text):,} chars from {file_path}")
        except Exception as e:
            logger.error(f"PyPDF2 also failed for {file_path}: {e}")

        return FileParser._clean_extracted_text(text)

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    @staticmethod
    def _parse_docx(file_path):
        """
        Extract text from a DOCX file preserving document order.

        Steps:
        1. First section's header and footer (common document metadata).
        2. Body via XML element-tree traversal so paragraphs and tables are
           interleaved in the order they appear in the document — unlike the
           python-docx default of iterating doc.paragraphs then doc.tables
           separately, which loses their relative position.

        Table cells are joined with a tab character; rows with a newline.
        """
        doc = docx.Document(file_path)
        parts = []

        # 1. Header / footer of the first section
        try:
            section = doc.sections[0]
            for hdr_ftr in (section.header, section.footer):
                if hdr_ftr:
                    for para in hdr_ftr.paragraphs:
                        t = para.text.strip()
                        if t:
                            parts.append(t)
        except Exception as hf_err:
            logger.debug(f"Could not extract headers/footers: {hf_err}")

        # 2. Body in document order
        for child in doc.element.body:
            tag = child.tag.split('}')[-1]  # strip XML namespace prefix

            if tag == 'p':
                # Collect all w:t text runs (handles hyperlinks, field codes, etc.)
                para_text = ''.join(
                    el.text for el in child.findall('.//' + qn('w:t'))
                    if el.text
                ).strip()

                # Collect any OMML equations embedded in this paragraph and
                # append them as [EQUATION: …] tags so downstream NLP can detect them.
                OMML_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/math'
                eq_tags = []
                for omath in child.findall(f'.//{{{OMML_NS}}}oMath'):
                    eq_text = FileParser._extract_omml_text(omath)
                    if eq_text:
                        eq_tags.append(f'[EQUATION: {eq_text}]')
                if eq_tags:
                    para_text = (para_text + ' ' if para_text else '') + ' '.join(eq_tags)

                if para_text:
                    parts.append(para_text)

            elif tag == 'tbl':
                # Rows → cells → paragraph text
                for row_el in child.findall('.//' + qn('w:tr')):
                    row_cells = []
                    for cell_el in row_el.findall(qn('w:tc')):
                        # All paragraphs within the cell
                        cell_parts = []
                        for para_el in cell_el.findall('.//' + qn('w:p')):
                            cell_text = ''.join(
                                el.text for el in para_el.findall('.//' + qn('w:t'))
                                if el.text
                            ).strip()
                            if cell_text:
                                cell_parts.append(cell_text)
                        if cell_parts:
                            row_cells.append(' '.join(cell_parts))
                    if row_cells:
                        parts.append('\t'.join(row_cells))

        text = '\n'.join(parts)
        logger.info(f"DOCX extracted {len(text):,} chars from {file_path}")
        return FileParser._clean_extracted_text(text)

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    @staticmethod
    def _shape_texts(shape):
        """
        Recursively extract text from a PPTX shape.

        Handles group shapes (which contain nested shapes) by recursing
        into their child collection rather than reading the concatenated
        group text (which would duplicate content).

        Args:
            shape: A python-pptx shape object.

        Returns:
            list[str]: Text strings found within the shape.
        """
        texts = []
        if hasattr(shape, 'shapes'):
            # GROUP shape — recurse into children
            for child_shape in shape.shapes:
                texts.extend(FileParser._shape_texts(child_shape))
        elif hasattr(shape, 'text') and shape.text.strip():
            texts.append(shape.text.strip())
        return texts

    @staticmethod
    def _parse_pptx(file_path):
        """
        Extract text from a PPTX file.

        For each slide:
        - All visible shape text (title, text boxes, placeholders, tables).
        - Grouped shapes are traversed recursively.
        - Speaker notes are appended with a [Notes] prefix.

        Deduplication within a single slide prevents the slide title
        appearing twice when it occurs in both the title placeholder and
        another shape.
        """
        prs = Presentation(file_path)
        slide_parts = []

        for idx, slide in enumerate(prs.slides, 1):
            parts = []
            seen = set()

            for shape in slide.shapes:
                for t in FileParser._shape_texts(shape):
                    if t not in seen:
                        seen.add(t)
                        parts.append(t)

            # Speaker notes
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f'[Notes] {notes_text}')
            except Exception:
                pass  # Not all slides have a notes frame

            if parts:
                slide_parts.append('\n'.join(parts))

        text = '\n\n'.join(slide_parts)
        logger.info(f"PPTX extracted {len(text):,} chars from {file_path}")
        return FileParser._clean_extracted_text(text)

    # ------------------------------------------------------------------
    # XLSX
    # ------------------------------------------------------------------

    @staticmethod
    def _parse_xlsx(file_path):
        """Extract all cell values from all sheets as tab-separated rows."""
        wb = load_workbook(file_path, data_only=True)
        parts = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append('\t'.join(cells))
        return FileParser._clean_extracted_text('\n'.join(parts))

    # ------------------------------------------------------------------
    # TXT
    # ------------------------------------------------------------------

    # ------------------------------------------------------------------
    # Image extraction  (PDF → PyMuPDF/fitz; DOCX → rels; PPTX → shapes)
    # ------------------------------------------------------------------

    @staticmethod
    def extract_images(file_path, file_type, output_dir):
        """
        Extract embedded images from a document and save them to *output_dir*.

        Args:
            file_path  (str): Absolute path to the source document.
            file_type  (str): Extension without the dot ('pdf', 'docx', 'pptx', …).
            output_dir (str): Directory where image files will be written.

        Returns:
            list[dict]: Each dict has keys:
                path        – absolute path of the saved image file
                page_number – PDF page / PPTX slide number (1-based), or None
                image_index – 0-based sequential index across the whole document
                width       – pixel width  (may be None if unavailable)
                height      – pixel height (may be None if unavailable)
        """
        ft = file_type.lower()
        if ft == 'pdf':
            return FileParser._extract_images_pdf(file_path, output_dir)
        elif ft in ('doc', 'docx'):
            return FileParser._extract_images_docx(file_path, output_dir)
        elif ft in ('ppt', 'pptx'):
            return FileParser._extract_images_pptx(file_path, output_dir)
        return []

    @staticmethod
    def _extract_images_pdf(file_path, output_dir):
        """Extract images from a PDF using PyMuPDF (fitz). Skipped gracefully if not installed."""
        images = []
        try:
            import fitz  # PyMuPDF
            os.makedirs(output_dir, exist_ok=True)
            doc = fitz.open(file_path)
            img_index = 0
            for page_num, page in enumerate(doc, 1):
                for img_info in page.get_images(full=True):
                    if img_index >= 50:
                        break
                    xref = img_info[0]
                    try:
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n > 4:          # CMYK → RGB
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        img_filename = f"page{page_num}_img{img_index}.png"
                        img_path = os.path.join(output_dir, img_filename)
                        pix.save(img_path)
                        images.append({
                            'path':        img_path,
                            'page_number': page_num,
                            'image_index': img_index,
                            'width':       pix.width,
                            'height':      pix.height,
                        })
                        img_index += 1
                    except Exception as img_err:
                        logger.debug(f"Skipped PDF image xref={xref}: {img_err}")
                if img_index >= 50:
                    break
            doc.close()
            logger.info(f"PDF image extraction: {len(images)} images saved to {output_dir}")
        except ImportError:
            logger.warning("PyMuPDF (fitz) not installed — PDF image extraction skipped. "
                           "Install with: pip install pymupdf")
        except Exception as e:
            logger.error(f"PDF image extraction error: {e}", exc_info=True)
        return images

    @staticmethod
    def _extract_images_docx(file_path, output_dir):
        """Extract embedded images from a DOCX file via its OPC relationships."""
        images = []
        try:
            import docx
            os.makedirs(output_dir, exist_ok=True)
            doc = docx.Document(file_path)
            img_index = 0
            for rel in doc.part.rels.values():
                if img_index >= 50:
                    break
                if 'image' not in rel.reltype:
                    continue
                try:
                    image_part = rel.target_part
                    img_data = image_part.blob
                    # Derive extension from content-type (e.g. image/png → png)
                    content_type = image_part.content_type
                    ext = content_type.split('/')[-1].replace('jpeg', 'jpg').split('+')[0]
                    ext = ext if ext.isalpha() else 'png'
                    img_filename = f"img{img_index}.{ext}"
                    img_path = os.path.join(output_dir, img_filename)
                    with open(img_path, 'wb') as f:
                        f.write(img_data)
                    images.append({
                        'path':        img_path,
                        'page_number': None,
                        'image_index': img_index,
                        'width':       None,
                        'height':      None,
                    })
                    img_index += 1
                except Exception as img_err:
                    logger.debug(f"Skipped DOCX image rel: {img_err}")
            logger.info(f"DOCX image extraction: {len(images)} images saved to {output_dir}")
        except Exception as e:
            logger.error(f"DOCX image extraction error: {e}", exc_info=True)
        return images

    @staticmethod
    def _extract_images_pptx(file_path, output_dir):
        """Extract picture shapes from a PPTX file."""
        images = []
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            os.makedirs(output_dir, exist_ok=True)
            prs = Presentation(file_path)
            img_index = 0
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if img_index >= 50:
                        break
                    # MSO_SHAPE_TYPE.PICTURE == 13
                    if shape.shape_type != 13:
                        continue
                    try:
                        image = shape.image
                        ext = image.ext or 'png'
                        img_filename = f"slide{slide_num}_img{img_index}.{ext}"
                        img_path = os.path.join(output_dir, img_filename)
                        with open(img_path, 'wb') as f:
                            f.write(image.blob)
                        images.append({
                            'path':        img_path,
                            'page_number': slide_num,
                            'image_index': img_index,
                            'width':       shape.width,
                            'height':      shape.height,
                        })
                        img_index += 1
                    except Exception as img_err:
                        logger.debug(f"Skipped PPTX image on slide {slide_num}: {img_err}")
                if img_index >= 50:
                    break
            logger.info(f"PPTX image extraction: {len(images)} images saved to {output_dir}")
        except Exception as e:
            logger.error(f"PPTX image extraction error: {e}", exc_info=True)
        return images

    # ------------------------------------------------------------------
    # OMML (Office Math Markup Language) helper
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_omml_text(element):
        """
        Return a plain-text representation of an OMML <m:oMath> element by
        concatenating all <m:t> text nodes separated by spaces.
        This gives a readable (though not perfectly rendered) equation string.
        """
        OMML_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/math'
        parts = []
        for el in element.iter(f'{{{OMML_NS}}}t'):
            if el.text and el.text.strip():
                parts.append(el.text.strip())
        return ' '.join(parts)

    @staticmethod
    def _parse_txt(file_path):
        """
        Read a plain-text file trying multiple encodings in priority order:
        UTF-8-sig (handles BOM), UTF-8, Latin-1, CP1252.

        If all fail, reads as raw bytes and decodes with errors='ignore'.
        """
        for enc in ('utf-8-sig', 'utf-8', 'latin-1', 'cp1252'):
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    text = f.read()
                logger.info(f"TXT read with encoding={enc} from {file_path}")
                return FileParser._clean_extracted_text(text)
            except (UnicodeDecodeError, LookupError):
                continue

        # Last resort: binary read, ignore undecodable bytes
        with open(file_path, 'rb') as f:
            text = f.read().decode('utf-8', errors='ignore')
        logger.warning(f"TXT read with errors='ignore' for {file_path}")
        return FileParser._clean_extracted_text(text)

    # ------------------------------------------------------------------
    # Image OCR (pytesseract)
    # ------------------------------------------------------------------

    @staticmethod
    def _parse_image(file_path):
        """
        Extract text from an image file using OCR (pytesseract).
        Requires: pytesseract + Tesseract engine installed on the system.
        Windows installer: https://github.com/UB-Mannheim/tesseract/wiki

        Returns cleaned OCR text, or empty string if OCR is unavailable.
        """
        if not HAS_PYTESSERACT:
            logger.warning(f"OCR unavailable (pytesseract not installed) for {file_path}")
            return ''
        try:
            from PIL import Image as PILImage
            img = PILImage.open(file_path).convert('RGB')
            text = pytesseract.image_to_string(img)
            logger.info(f"OCR extracted {len(text.split())} words from {file_path}")
            return FileParser._clean_extracted_text(text)
        except Exception as e:
            logger.error(f"OCR failed for {file_path}: {e}", exc_info=True)
            return ''
